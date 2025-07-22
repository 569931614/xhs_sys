#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import json
import logging
import tempfile
import requests
import uuid
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
import shutil
import subprocess
import time
import random

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),  # 输出到控制台
    ]
)

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

def download_image(image_url, output_path):
    """
    下载图片并保存到指定路径
    """
    try:
        response = requests.get(image_url, stream=True)
        response.raise_for_status()
        
        with open(output_path, 'wb') as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)
                
        # 验证图片是否正确下载
        if os.path.getsize(output_path) == 0:
            logger.error(f"下载的图片为空: {image_url}")
            return False
                
        logger.info(f"成功下载图片: {image_url} -> {output_path}")
        return True
    except Exception as e:
        logger.error(f"下载图片失败: {image_url}, 错误: {str(e)}")
        return False

def replace_ppt_images_and_text(template_path, images, slogans, output_path, textbox_mapping=None, use_direct_replacement=True):
    """
    替换PPT中的图片和文字
    """
    # 生成唯一的临时目录ID，防止并发冲突
    temp_dir_id = str(uuid.uuid4())
    temp_image_paths = []
    task_temp_dir = None
    
    try:
        # 确保至少有一个图片或标语
        if not images and not slogans:
            logger.error("图片和标语不能同时为空")
            return False
            
        # 创建唯一的临时目录
        base_temp_dir = tempfile.gettempdir()
        task_temp_dir = os.path.join(base_temp_dir, temp_dir_id)
        os.makedirs(task_temp_dir, exist_ok=True)
        logger.info(f"为任务创建临时目录: {task_temp_dir}")
            
        # 下载图片
        image_paths = []
        if images:
            # 如果图片数量超过PPT中的图片数量，随机选取
            ppt = Presentation(template_path)
            picture_count = sum(1 for slide in ppt.slides for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)
            logger.info(f"PPT中的图片数量: {picture_count}")
            
            if len(images) > picture_count:
                logger.info(f"上传的图片数量({len(images)})多于PPT中的图片数量({picture_count})，将随机选取")
                selected_images = random.sample(images, picture_count)
                logger.info(f"随机选取的图片: {selected_images}")
            else:
                selected_images = images
                
            for i, image_url in enumerate(selected_images):
                # 使用任务特定的临时目录存放图片
                image_path = os.path.join(task_temp_dir, f"image_{i}.jpg")
                temp_image_paths.append(image_path)  # 记录需要清理的图片路径
                
                # 下载图片
                if not download_image(image_url, image_path):
                    continue
                
                image_paths.append(image_path)
        
        # 复制模板到输出路径
        output_dir = os.path.dirname(output_path)
        if not os.path.exists(output_dir):
            os.makedirs(output_dir, exist_ok=True)
            
        shutil.copy2(template_path, output_path)
        logger.info(f"已复制模板到输出路径: {output_path}")
        
        # 打开PPT
        prs = Presentation(output_path)
        
        # 处理所有页面
        for slide_index, slide in enumerate(prs.slides):
            logger.info(f"处理第 {slide_index+1} 页")
            
            # 替换图片
            if image_paths:
                # 获取所有形状
                shapes = slide.shapes
                
                # 查找图片占位符
                picture_placeholders = []
                for shape in shapes:
                    if hasattr(shape, 'shape_type') and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        picture_placeholders.append(shape)
                
                # 替换图片
                for i, (placeholder, image_path) in enumerate(zip(picture_placeholders, image_paths)):
                    try:
                        # 验证图片文件是否存在
                        if not os.path.exists(image_path):
                            logger.error(f"图片文件不存在: {image_path}")
                            continue
                            
                        # 获取占位符的位置和大小
                        left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
                        
                        # 删除占位符
                        placeholder._element.getparent().remove(placeholder._element)
                        
                        # 添加新图片
                        new_picture = slide.shapes.add_picture(image_path, left, top, width, height)
                        
                        # 将图片移动到最底层
                        shapes_tree = slide.shapes._spTree
                        picture_element = new_picture._element
                        # 从当前位置删除
                        shapes_tree.remove(picture_element)
                        # 在最前面插入，使其成为最底层元素
                        shapes_tree.insert(0, picture_element)
                        
                        logger.info(f"已替换第 {slide_index+1} 页图片 {i+1}/{len(image_paths)}，并将其设置为最底层")
                    except Exception as e:
                        logger.error(f"替换图片失败: {str(e)}")
                        return False
            
            # 替换文字 - 采用更直接的方法
            if slogans and len(slogans) > 0:
                # 只处理有文本框的幻灯片
                text_boxes = []
                for shape in slide.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text:
                        text_boxes.append(shape)
                        logger.info(f"找到文本框: {shape.text}")
                
                logger.info(f"本页共找到 {len(text_boxes)} 个非空文本框")
                logger.info(f"可用标语列表: {slogans}")
                
                # 如果找到了文本框，就按顺序替换
                for i, text_box in enumerate(text_boxes):
                    if i < len(slogans):
                        # 完全不考虑原文，直接用新文字替换
                        slogan_to_use = slogans[i]
                        logger.info(f"准备替换第 {i+1} 个文本框: 原文='{text_box.text}', 新文='{slogan_to_use}'")
                        
                        try:
                            # 保存原有段落样式
                            original_paragraph = text_box.text_frame.paragraphs[0]
                            original_alignment = original_paragraph.alignment
                            
                            # 获取第一个运行对象的样式
                            if original_paragraph.runs:
                                first_run = original_paragraph.runs[0]
                                font_name = first_run.font.name
                                font_size = first_run.font.size
                                font_bold = first_run.font.bold
                                font_italic = first_run.font.italic
                                font_underline = first_run.font.underline
                                
                                # 保存颜色信息
                                original_color = None
                                try:
                                    if hasattr(first_run.font.color, '_element'):
                                        original_color = first_run.font.color._element
                                except Exception as e:
                                    logger.warning(f"获取字体颜色失败: {str(e)}")
                            
                            # 清除原有文本
                            text_box.text_frame.clear()
                            
                            # 添加新文本
                            new_paragraph = text_box.text_frame.paragraphs[0] if text_box.text_frame.paragraphs else text_box.text_frame.add_paragraph()
                            new_run = new_paragraph.add_run()
                            new_run.text = slogan_to_use
                            
                            # 应用原有样式
                            if 'first_run' in locals():
                                try:
                                    new_run.font.name = font_name
                                    new_run.font.size = font_size
                                    new_run.font.bold = font_bold
                                    new_run.font.italic = font_italic
                                    new_run.font.underline = font_underline
                                    
                                    # 应用颜色
                                    if original_color is not None:
                                        new_run.font.color._element = original_color
                                except Exception as e:
                                    logger.warning(f"应用样式失败: {str(e)}")
                            
                            # 设置段落对齐方式
                            if 'original_alignment' in locals():
                                new_paragraph.alignment = original_alignment
                            
                            # 将文本框移到最上层
                            try:
                                spTree = slide.shapes._spTree
                                spTree.remove(text_box._element)
                                spTree.append(text_box._element)
                            except Exception as e:
                                logger.warning(f"移动文本框到最上层失败: {str(e)}")
                            
                            logger.info(f"成功替换文本框: 新文本='{slogan_to_use}'")
                        except Exception as e:
                            logger.error(f"替换文本失败: {str(e)}")
                    else:
                        logger.warning(f"标语数量不足，无法替换第 {i+1} 个文本框")
            
            # 如果使用了文本映射方式，也尝试按映射替换
            elif textbox_mapping:
                # 获取所有形状
                shapes = slide.shapes
                
                # 查找文本框
                text_boxes = []
                for shape in shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        text_boxes.append(shape)
                        logger.info(f"找到文本框: {shape.text if shape.text else '空文本框'}")
                
                # 添加调试信息
                logger.info(f"文本映射数据: {textbox_mapping}")
                
                # 替换文字 - 遍历所有文本框
                for text_box in text_boxes:
                    try:
                        # 检查文本框是否有文本
                        if not text_box.text:
                            continue
                            
                        # 遍历文本映射
                        for slogan, text_to_replace in textbox_mapping.items():
                            # 添加更详细的日志，包括类型信息和字符编码
                            logger.info(f"检查替换: 新文本='{slogan}'({type(slogan)}), 要替换的文本='{text_to_replace}'({type(text_to_replace)})")
                            logger.info(f"文本框内容: '{text_box.text}'({type(text_box.text)})")
                            
                            # 将所有输入标准化为字符串
                            if not isinstance(text_to_replace, str):
                                text_to_replace = str(text_to_replace)
                            if not isinstance(text_box.text, str):
                                box_text = str(text_box.text)
                            else:
                                box_text = text_box.text
                            
                            # 尝试多种匹配方式
                            found_match = False
                            # 1. 精确匹配
                            if text_to_replace and text_to_replace in box_text:
                                logger.info(f"精确匹配成功: '{text_to_replace}' 在 '{box_text}' 中")
                                found_match = True
                            # 2. 忽略大小写的匹配
                            elif text_to_replace and text_to_replace.lower() in box_text.lower():
                                logger.info(f"忽略大小写匹配成功: '{text_to_replace.lower()}' 在 '{box_text.lower()}' 中")
                                found_match = True
                            # 3. 移除空白字符后的匹配
                            elif text_to_replace and text_to_replace.strip() and text_to_replace.strip() in box_text.strip():
                                logger.info(f"移除空白字符后匹配成功: '{text_to_replace.strip()}' 在 '{box_text.strip()}' 中")
                                found_match = True
                                
                            if not found_match:
                                logger.warning(f"没有找到匹配项，跳过此文本框")
                                continue
                                
                            # 找到匹配的文本，进行替换
                            logger.info(f"找到匹配的文本! 将替换为 '{slogan}'")
                            
                            # 保存原有段落样式
                            original_paragraph = text_box.text_frame.paragraphs[0]
                            original_alignment = original_paragraph.alignment
                            
                            # 获取第一个运行对象的样式作为默认样式
                            if original_paragraph.runs:
                                first_run = original_paragraph.runs[0]
                                font_name = first_run.font.name
                                font_size = first_run.font.size
                                font_bold = first_run.font.bold
                                font_italic = first_run.font.italic
                                font_underline = first_run.font.underline
                                
                                # 保存原始颜色信息
                                original_color = None
                                try:
                                    if hasattr(first_run.font.color, '_element'):
                                        original_color = first_run.font.color._element
                                        logger.info("成功获取原始颜色元素")
                                    else:
                                        logger.warning("原始文本没有颜色元素")
                                except Exception as e:
                                    logger.warning(f"获取原始颜色信息失败: {str(e)}")
                                
                                # 清除原有文本
                                text_box.text_frame.clear()
                                
                                # 添加新段落并应用样式
                                new_paragraph = text_box.text_frame.paragraphs[0] if text_box.text_frame.paragraphs else text_box.text_frame.add_paragraph()
                                new_run = new_paragraph.add_run()
                                new_run.text = slogan  # 直接使用新文本
                                
                                # 应用原有样式
                                if original_paragraph.runs:
                                    new_run.font.name = font_name
                                    new_run.font.size = font_size
                                    new_run.font.bold = font_bold
                                    new_run.font.italic = font_italic
                                    new_run.font.underline = font_underline
                                    
                                    # 应用颜色
                                    try:
                                        if original_color is not None:
                                            # 直接使用原始颜色元素
                                            new_run.font.color._element = original_color
                                            logger.info("成功应用原始颜色元素")
                                        else:
                                            # 如果无法获取原始颜色，尝试设置为主题白色
                                            new_run.font.color.theme_color = MSO_THEME_COLOR.BACKGROUND_1  # 背景1通常是白色
                                            logger.info("使用主题白色作为替代")
                                    except Exception as e:
                                        logger.warning(f"设置字体颜色失败: {str(e)}")
                                        try:
                                            # 最后尝试使用RGB白色
                                            new_run.font.color.rgb = 16777215  # 白色RGB值
                                            logger.info("使用RGB白色作为最后替代")
                                        except Exception as e:
                                            logger.error(f"所有颜色设置方法都失败: {str(e)}")
                                
                                # 设置段落对齐方式
                                new_paragraph.alignment = original_alignment
                                
                                # 确保文本框在最上层
                                spTree = slide.shapes._spTree
                                spTree.remove(text_box._element)
                                spTree.append(text_box._element)
                                
                                logger.info(f"已替换文字: {text_to_replace} -> {slogan}")
                                break
                                
                    except Exception as e:
                        logger.error(f"替换文字失败: {str(e)}")
                        return False
        
        # 保存修改后的PPT - 替换为重试版本
        try:
            save_ppt_with_retry(prs, output_path)
            logger.info(f"已保存修改后的PPT: {output_path}")
        except Exception as e:
            logger.error(f"替换PPT内容失败: {str(e)}")
            return False
        
        return True
        
    except Exception as e:
        logger.error(f"替换PPT内容失败: {str(e)}")
        return False
    finally:
        # 清理临时文件
        try:
            # 清理下载的图片
            for path in temp_image_paths:
                try:
                    if os.path.exists(path):
                        os.remove(path)
                        logger.info(f"已清理临时图片: {path}")
                except Exception as cleanup_err:
                    logger.warning(f"清理临时图片失败: {path}, 错误: {str(cleanup_err)}")
            
            # 清理临时目录（如果存在且为空）
            if task_temp_dir and os.path.exists(task_temp_dir):
                try:
                    if not os.listdir(task_temp_dir):
                        os.rmdir(task_temp_dir)
                        logger.info(f"已清理空的临时目录: {task_temp_dir}")
                    else:
                        logger.warning(f"临时目录不为空，无法删除: {task_temp_dir}")
                except Exception as dir_err:
                    logger.warning(f"清理临时目录失败: {task_temp_dir}, 错误: {str(dir_err)}")
        except Exception as e:
            logger.warning(f"清理临时资源失败: {str(e)}")

def save_ppt_as_images(ppt_path, output_dir):
    """
    将PPT保存为图片
    """
    try:
        logger.info(f"开始将PPT保存为图片: {ppt_path}")
        logger.info(f"输出目录: {output_dir}")
        
        # 确保输出目录存在
        os.makedirs(output_dir, exist_ok=True)
        logger.info("输出目录已创建")
        
        # 打开PPT
        prs = Presentation(ppt_path)
        logger.info(f"成功打开PPT，共 {len(prs.slides)} 页")
        
        # 处理所有页面
        slide_count = len(prs.slides)
        if slide_count > 0:
            # 创建临时目录用于处理
            temp_dir = tempfile.mkdtemp()
            temp_output_path = os.path.join(temp_dir, "temp.pptx")
            prs.save(temp_output_path)
            logger.info(f"已保存临时PPT文件: {temp_output_path}")
            
            # 尝试使用libreoffice转换
            try:
                # 在Windows和Linux上的命令会有所不同
                if os.name == 'nt':  # Windows
                    try:
                        # 尝试使用PowerPoint自动化
                        import win32com.client
                        logger.info("尝试使用PowerPoint自动化转换PPT为图片")
                        
                        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                        powerpoint.Visible = True
                        
                        # 添加异常处理和详细日志
                        try:
                            presentation = powerpoint.Presentations.Open(temp_output_path)
                            logger.info(f"成功打开演示文稿，幻灯片数: {presentation.Slides.Count}")
                            
                            # 使用内置的导出功能导出整个演示文稿为PNG
                            try:
                                # 方法1: 尝试使用内置导出功能
                                export_path = os.path.join(output_dir, "export")
                                logger.info(f"尝试导出整个演示文稿到: {export_path}")
                                presentation.ExportAsFixedFormat(export_path, 2, PrintRange=1)  # 2 = ppFixedFormatTypePNG
                                logger.info("成功导出演示文稿")
                            except Exception as export_err:
                                logger.warning(f"导出整个演示文稿失败: {str(export_err)}，尝试单张导出")
                                
                                # 方法2: 逐页导出
                                for i in range(1, presentation.Slides.Count + 1):
                                    try:
                                        slide_path = os.path.join(output_dir, f"slide{i}.png")
                                        logger.info(f"尝试导出第 {i} 页: {slide_path}")
                                        # 使用正确的索引访问方式，PowerPoint COM从1开始索引
                                        presentation.Slides.Item(i).Export(slide_path, "PNG")
                                        logger.info(f"已导出第 {i} 页: {slide_path}")
                                    except Exception as slide_err:
                                        logger.error(f"导出第 {i} 页失败: {str(slide_err)}")
                        except Exception as ppt_err:
                            logger.error(f"PowerPoint操作失败: {str(ppt_err)}")
                            raise
                        finally:
                            # 确保关闭PowerPoint
                            try:
                                if 'presentation' in locals():
                                    presentation.Close()
                                powerpoint.Quit()
                                logger.info("已关闭PowerPoint")
                            except Exception as close_err:
                                logger.warning(f"关闭PowerPoint失败: {str(close_err)}")
                        
                    except Exception as win_err:
                        logger.error(f"使用PowerPoint自动化转换失败: {str(win_err)}")
                        # 如果PowerPoint自动化失败，尝试使用备用方法
                        self_convert_ppt_to_images(prs, output_dir)
                else:  # Linux 或 macOS
                    # 使用libreoffice转换
                    cmd = [
                        "libreoffice", "--headless", "--convert-to", "png",
                        "--outdir", output_dir, temp_output_path
                    ]
                    logger.info(f"执行命令: {' '.join(cmd)}")
                    
                    process = subprocess.Popen(
                        cmd,
                        stdout=subprocess.PIPE,
                        stderr=subprocess.PIPE,
                        text=True
                    )
                    
                    stdout, stderr = process.communicate(timeout=120)
                    
                    if process.returncode != 0:
                        logger.error(f"LibreOffice转换失败: {stderr}")
                        # 如果libreoffice转换失败，尝试使用备用方法
                        self_convert_ppt_to_images(prs, output_dir)
                    
                    logger.info(f"LibreOffice转换输出: {stdout}")
                    
                # 检查是否生成了任何图片文件
                files = os.listdir(output_dir)
                png_files = [f for f in files if f.endswith('.png')]
                
                if not png_files:
                    logger.warning("未检测到生成的图片文件，尝试使用备用方法")
                    self_convert_ppt_to_images(prs, output_dir)
                    
                    # 再次检查图片文件
                    files = os.listdir(output_dir)
                    png_files = [f for f in files if f.endswith('.png')]
                
                # 重命名生成的图片以保持一致的命名格式
                png_files.sort()
                
                for i, file in enumerate(png_files, 1):
                    old_path = os.path.join(output_dir, file)
                    new_path = os.path.join(output_dir, f"slide{i}.png")
                    if old_path != new_path:
                        os.rename(old_path, new_path)
                        logger.info(f"已重命名图片: {old_path} -> {new_path}")
                
                logger.info(f"共生成 {len(png_files)} 张图片")
                return True
                
            except Exception as e:
                logger.error(f"转换PPT为图片失败: {str(e)}")
                # 使用备用方法
                return self_convert_ppt_to_images(prs, output_dir)
            finally:
                # 清理临时目录
                try:
                    shutil.rmtree(temp_dir)
                    logger.info(f"已清理临时目录: {temp_dir}")
                except Exception as cleanup_err:
                    logger.warning(f"清理临时目录失败: {temp_dir}, 错误: {str(cleanup_err)}")
        else:
            logger.error("PPT中没有幻灯片")
            return False
            
    except Exception as e:
        logger.error(f"保存PPT为图片时发生错误: {str(e)}")
        return False

def self_convert_ppt_to_images(prs, output_dir):
    """
    不依赖外部程序，使用python-pptx直接将PPT转换为图片
    这是一个备用方法，当PowerPoint自动化或LibreOffice转换失败时使用
    
    Args:
        prs: PowerPoint演示文稿对象
        output_dir: 输出目录
        
    Returns:
        bool: 是否成功转换
    """
    try:
        logger.info("使用备用方法转换PPT为图片")
        
        # 使用简单的方法 - 为每个幻灯片创建一个文本表示
        for i, slide in enumerate(prs.slides, 1):
            # 创建一个文本文件作为占位符
            slide_path = os.path.join(output_dir, f"slide{i}.png")
            
            # 创建一个简单的文本图像，说明需要手动查看PPT
            try:
                from PIL import Image, ImageDraw, ImageFont
                
                # 创建一个白色背景图像
                img = Image.new('RGB', (800, 600), color=(255, 255, 255))
                d = ImageDraw.Draw(img)
                
                # 尝试加载字体，如果失败则使用默认字体
                try:
                    font = ImageFont.truetype("arial.ttf", 20)
                except Exception:
                    font = ImageFont.load_default()
                
                # 绘制文本
                text = f"幻灯片 {i} (请查看原始PPT文件获取完整内容)"
                d.text((50, 50), text, fill=(0, 0, 0), font=font)
                
                # 添加幻灯片上的文本内容
                text_content = []
                y_pos = 100
                
                for shape in slide.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text_line = "".join(run.text for run in paragraph.runs)
                            if text_line.strip():
                                text_content.append(text_line)
                                d.text((50, y_pos), text_line[:60], fill=(0, 0, 0), font=font)
                                y_pos += 30
                
                # 保存图像
                img.save(slide_path)
                logger.info(f"已创建幻灯片预览图: {slide_path}")
                
            except ImportError:
                # 如果PIL不可用，创建空文件
                with open(slide_path, 'wb') as f:
                    f.write(b'')
                logger.warning(f"PIL不可用，已创建空文件: {slide_path}")
        
        logger.info(f"已使用备用方法处理 {len(prs.slides)} 张幻灯片")
        return True
        
    except Exception as e:
        logger.error(f"备用方法转换PPT失败: {str(e)}")
        return False

# 添加一个使用临时文件策略保存PPT的函数
def save_ppt_with_retry(prs, output_path, max_retries=3, delay=1):
    """
    使用临时文件策略和重试机制保存PPT文件
    
    Args:
        prs: PowerPoint演示文稿对象
        output_path: 输出路径
        max_retries: 最大重试次数
        delay: 重试间隔(秒)
    
    Returns:
        bool: 是否成功保存
    """
    # 创建临时文件名
    temp_dir = os.path.dirname(output_path)
    temp_filename = f"temp_{uuid.uuid4().hex}.pptx"
    temp_path = os.path.join(temp_dir, temp_filename)
    
    for attempt in range(max_retries):
        try:
            # 先保存到临时文件
            logger.info(f"尝试保存到临时文件: {temp_path}")
            prs.save(temp_path)
            
            # 直接使用复制方式而不是重命名
            # 这样可以避免目标文件存在时的冲突
            try:
                logger.info(f"尝试将临时文件复制到目标位置: {output_path}")
                shutil.copy2(temp_path, output_path)
                logger.info(f"第 {attempt+1} 次尝试成功保存修改后的PPT: {output_path}")
                
                # 尝试删除临时文件
                try:
                    os.remove(temp_path)
                    logger.info(f"已删除临时文件: {temp_path}")
                except Exception as rm_err:
                    logger.warning(f"删除临时文件失败: {str(rm_err)}")
                
                return True
            except Exception as copy_err:
                logger.error(f"复制文件失败: {str(copy_err)}")
                
                # 如果复制失败，尝试直接返回临时文件
                if os.path.exists(temp_path):
                    logger.info(f"使用临时文件作为输出: {temp_path} -> {output_path}")
                    # 创建一个新的输出路径，不与原路径冲突
                    new_output = f"{os.path.splitext(output_path)[0]}_{uuid.uuid4().hex[:8]}.pptx"
                    os.rename(temp_path, new_output)
                    logger.info(f"已重命名临时文件为新输出: {new_output}")
                    # 将新路径添加到日志中，后续可以从中提取
                    logger.info(f"ALTERNATE_OUTPUT_PATH:{new_output}")
                    return True
                raise
                
        except PermissionError as e:
            logger.warning(f"保存PPT文件权限错误 (尝试 {attempt+1}/{max_retries}): {str(e)}")
            if attempt < max_retries - 1:
                time.sleep(delay)  # 等待一段时间后重试
            else:
                logger.error(f"达到最大重试次数，无法保存PPT文件: {str(e)}")
                
                # 如果临时文件存在，尝试使用临时文件
                if os.path.exists(temp_path):
                    try:
                        # 创建一个新的输出路径，不与原路径冲突
                        new_output = f"{os.path.splitext(output_path)[0]}_{uuid.uuid4().hex[:8]}.pptx"
                        os.rename(temp_path, new_output)
                        logger.info(f"已重命名临时文件为新输出: {new_output}")
                        # 将新路径添加到日志中，后续可以从中提取
                        logger.info(f"ALTERNATE_OUTPUT_PATH:{new_output}")
                        return True
                    except Exception as rename_err:
                        logger.error(f"重命名临时文件失败: {str(rename_err)}")
                
                raise
        except Exception as e:
            logger.error(f"保存PPT文件时发生其他错误: {str(e)}")
            # 尝试清理临时文件
            try:
                if os.path.exists(temp_path):
                    os.remove(temp_path)
                    logger.info(f"已删除临时文件: {temp_path}")
            except Exception as rm_err:
                logger.warning(f"删除临时文件失败: {str(rm_err)}")
            raise
            
    # 尝试清理临时文件
    try:
        if os.path.exists(temp_path):
            os.remove(temp_path)
            logger.info(f"已删除临时文件: {temp_path}")
    except Exception as rm_err:
        logger.warning(f"删除临时文件失败: {str(rm_err)}")
        
    return False

def main():
    """
    主函数
    """
    if len(sys.argv) < 2:
        print("用法: python ppt_processor.py <参数JSON文件> [输出目录]")
        sys.exit(1)
        
    params_file = sys.argv[1]
    output_image_dir = sys.argv[2] if len(sys.argv) > 2 else None
    
    try:
        with open(params_file, 'r', encoding='utf-8') as f:
            params = json.load(f)
            
        # 解析参数
        template_path = params.get('template_path')
        output_path = params.get('output_path')
        images = params.get('images', [])
        slogans = params.get('slogans', [])
        textbox_mapping = params.get('textbox_mapping', {})
        use_direct_replacement = params.get('use_direct_replacement', True)  # 新增选项，默认为True
        
        logger.info(f"收到的标语列表: {slogans}")
        logger.info(f"文本映射: {textbox_mapping}")
        logger.info(f"是否使用直接替换: {use_direct_replacement}")
        
        # PPT处理
        if replace_ppt_images_and_text(template_path, images, slogans, output_path, textbox_mapping, use_direct_replacement):
            logger.info("PPT内容替换成功")
            
            # 如果指定了输出图片目录，转换PPT为图片
            if output_image_dir:
                if save_ppt_as_images(output_path, output_image_dir):
                    logger.info("PPT转换为图片成功")
                    print(json.dumps({"success": True, "message": "PPT处理成功"}))
                    sys.exit(0)
                else:
                    logger.error("PPT转换为图片失败")
                    print(json.dumps({"success": False, "message": "PPT转换为图片失败"}))
                    sys.exit(1)
            else:
                print(json.dumps({"success": True, "message": "PPT内容替换成功"}))
                sys.exit(0)
        else:
            logger.error("PPT内容替换失败")
            print(json.dumps({"success": False, "message": "PPT内容替换失败"}))
            sys.exit(1)
            
    except Exception as e:
        logger.error(f"处理PPT时发生错误: {str(e)}")
        print(json.dumps({"success": False, "message": f"处理PPT时发生错误: {str(e)}"}))
        sys.exit(1)

if __name__ == "__main__":
    main() 