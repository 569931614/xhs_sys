#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import sys
import json
import logging
import tempfile
import uuid
from pptx import Presentation
import shutil
import subprocess
import time

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.StreamHandler(),  # 输出到控制台
    ]
)

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)

def save_ppt_as_images(ppt_path, output_dir):
    """
    将PPT保存为图片
    """
    try:
        logger.info(f"开始将PPT保存为图片: {ppt_path}")
        logger.info(f"输出目录: {output_dir}")
        
        # 确保输出目录存在
        os.makedirs(output_dir, exist_ok=True)
        logger.info("输出目录已创建")
        
        # 打开PPT
        prs = Presentation(ppt_path)
        logger.info(f"成功打开PPT，共 {len(prs.slides)} 页")
        
        # 处理所有页面
        slide_count = len(prs.slides)
        if slide_count > 0:
            # 创建临时目录用于处理
            temp_dir = tempfile.mkdtemp()
            temp_output_path = os.path.join(temp_dir, "temp.pptx")
            prs.save(temp_output_path)
            logger.info(f"已保存临时PPT文件: {temp_output_path}")
            
            # 尝试使用libreoffice转换
            try:
                # 在Windows和Linux上的命令会有所不同
                if os.name == 'nt':  # Windows
                    try:
                        # 尝试使用PowerPoint自动化
                        import win32com.client
                        logger.info("尝试使用PowerPoint自动化转换PPT为图片")
                        
                        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                        powerpoint.Visible = True
                        presentation = powerpoint.Presentations.Open(temp_output_path)
                        
                        for i in range(1, slide_count + 1):
                            slide_path = os.path.join(output_dir, f"slide{i}.png")
                            presentation.Slides[i].Export(slide_path, "PNG")
                            logger.info(f"已导出第 {i} 页: {slide_path}")
                            
                        presentation.Close()
                        powerpoint.Quit()
                    except Exception as win_err:
                        logger.error(f"使用PowerPoint自动化转换失败: {str(win_err)}")
                        raise
                else:  # Linux 或 macOS
                    # 使用libreoffice转换
                    cmd = [
                        "libreoffice", "--headless", "--convert-to", "png",
                        "--outdir", output_dir, temp_output_path
                    ]
                    logger.info(f"执行命令: {' '.join(cmd)}")
                    
                    process = subprocess.Popen(
                        cmd,
                        stdout=subprocess.PIPE,
                        stderr=subprocess.PIPE,
                        text=True
                    )
                    
                    stdout, stderr = process.communicate(timeout=120)
                    
                    if process.returncode != 0:
                        logger.error(f"LibreOffice转换失败: {stderr}")
                        raise Exception(f"LibreOffice转换失败: {stderr}")
                    
                    logger.info(f"LibreOffice转换输出: {stdout}")
                    
                # 重命名生成的图片以保持一致的命名格式
                files = os.listdir(output_dir)
                png_files = [f for f in files if f.endswith('.png')]
                png_files.sort()
                
                for i, file in enumerate(png_files, 1):
                    old_path = os.path.join(output_dir, file)
                    new_path = os.path.join(output_dir, f"slide{i}.png")
                    if old_path != new_path:
                        os.rename(old_path, new_path)
                        logger.info(f"已重命名图片: {old_path} -> {new_path}")
                
                logger.info(f"共生成 {len(png_files)} 张图片")
                return True
                
            except Exception as e:
                logger.error(f"转换PPT为图片失败: {str(e)}")
                return False
            finally:
                # 清理临时目录
                try:
                    shutil.rmtree(temp_dir)
                    logger.info(f"已清理临时目录: {temp_dir}")
                except Exception as cleanup_err:
                    logger.warning(f"清理临时目录失败: {temp_dir}, 错误: {str(cleanup_err)}")
        else:
            logger.error("PPT中没有幻灯片")
            return False
            
    except Exception as e:
        logger.error(f"保存PPT为图片时发生错误: {str(e)}")
        return False

def main():
    """
    主函数
    """
    if len(sys.argv) < 2:
        print("用法: python ppt_to_images.py <参数JSON文件>")
        sys.exit(1)
        
    params_file = sys.argv[1]
    
    try:
        with open(params_file, 'r', encoding='utf-8') as f:
            params = json.load(f)
            
        # 解析参数
        ppt_path = params.get('ppt_path')
        output_dir = params.get('output_dir')
        
        if not ppt_path or not output_dir:
            error_msg = "参数文件中缺少必要的参数: ppt_path 或 output_dir"
            logger.error(error_msg)
            print(json.dumps({"success": False, "message": error_msg}))
            sys.exit(1)
        
        # 转换PPT为图片
        if save_ppt_as_images(ppt_path, output_dir):
            logger.info("PPT转换为图片成功")
            print(json.dumps({"success": True, "message": "PPT转换为图片成功"}))
            sys.exit(0)
        else:
            logger.error("PPT转换为图片失败")
            print(json.dumps({"success": False, "message": "PPT转换为图片失败"}))
            sys.exit(1)
            
    except Exception as e:
        logger.error(f"处理PPT时发生错误: {str(e)}")
        print(json.dumps({"success": False, "message": f"处理PPT时发生错误: {str(e)}"}))
        sys.exit(1)

if __name__ == "__main__":
    main() 